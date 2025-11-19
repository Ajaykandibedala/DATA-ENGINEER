try:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore
except Exception:
    Presentation = None
    Inches = None
    Pt = None

from pathlib import Path


def main():
    output_dir = Path("output")
    output_dir.mkdir(parents=True, exist_ok=True)

    if Presentation is None:
        print("python-pptx is not installed — skipping report generation.")
        print("Install it with: pip install python-pptx")
        return

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # set title safely (some layouts may not have a title placeholder)
    title_shape = slide.shapes.title
    if title_shape is not None:
        title_shape.text = "EDA Summary: IndiaMART"
    else:
        tx_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        tf = tx_box.text_frame
        p = tf.paragraphs[0]
        p.text = "EDA Summary: IndiaMART"
        p.font.size = Pt(24)

    # add an image slide
    img_path = output_dir / "price_histogram.png"
    if img_path.exists():
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        slide2.shapes.add_picture(str(img_path), Inches(1), Inches(1.5), width=Inches(8))

    prs.save(str(output_dir / "report.pptx"))
    print(f"Saved {output_dir / 'report.pptx'}")


if __name__ == "__main__":
    main()
